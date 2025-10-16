from typing import Dict, List, Optional, Set

from pptx import Presentation
from tqdm import tqdm

from .glossary import apply_glossary, read_glossary
from .pptx_utils import (
    TextItem,
    collect_items,
    iter_text_frames,
    parse_slide_spec,
    set_text,
    summarize_deck,
)
from .translators import BaseTranslator


class TranslationResult:
    def __init__(self, skipped_all: bool, warnings: Optional[List[str]] = None):
        self.skipped_all = skipped_all
        self.warnings = warnings or []


def translate_presentation(
    translator: BaseTranslator,
    input_path: str,
    output_path: str,
    include_notes: bool,
    include_masters: bool,
    glossary_path: Optional[str],
    strategy: str,
    dry_run: bool,
    log_path: Optional[str],
    slide_spec: Optional[str],
) -> TranslationResult:
    prs = Presentation(input_path)
    items, slide_titles = collect_items(
        prs, include_notes=include_notes, include_masters=include_masters
    )

    selected_slides, slides_invalid = parse_slide_spec(slide_spec, len(prs.slides))
    warnings: List[str] = []
    if slides_invalid:
        warnings.append(
            "[WARN] Some slide identifiers in --slides were invalid or out of range and have been ignored."
        )
    if slide_spec and selected_slides is not None and not selected_slides:
        warnings.append(
            "[WARN] No valid slide numbers specified in --slides; nothing will be translated."
        )

    slide_indices: List[Optional[int]] = []
    for it in items:
        sidx = None
        if it.owner.startswith("slide["):
            try:
                sidx = int(it.owner.split("[")[1].split("]")[0]) - 1
            except Exception:
                sidx = None
        slide_indices.append(sidx)

    glossary = read_glossary(glossary_path)

    log_fh = None
    if log_path:
        log_fh = open(log_path, "w", encoding="utf-8")

    try:
        by_slide: Dict[Optional[int], List[TextItem]] = {}
        for it, sidx in zip(items, slide_indices):
            by_slide.setdefault(sidx, []).append(it)

        deck_summary = summarize_deck(slide_titles)
        progress = tqdm(total=len(items), desc="Translating")

        def make_context(current_idx: Optional[int]) -> str:
            if strategy == "title-only":
                title = (
                    slide_titles[current_idx]
                    if (current_idx is not None and 0 <= current_idx < len(slide_titles))
                    else ""
                )
                return f"Deck title(s): {slide_titles[0] if slide_titles else ''}. Current slide title: {title}."
            elif strategy == "deck":
                return deck_summary
            else:  # neighbor
                if current_idx is None:
                    return deck_summary
                titles = []
                for j in [current_idx - 1, current_idx, current_idx + 1]:
                    if 0 <= j < len(slide_titles):
                        titles.append(f"Slide {j+1}: {slide_titles[j]}")
                return "Neighbor titles:\n" + "\n".join(titles)

        def should_translate_group(group_slide_idx: Optional[int]) -> bool:
            if selected_slides is None:
                return True
            if group_slide_idx is None:
                return False
            return group_slide_idx in selected_slides

        any_translated = False

        for sidx, group in by_slide.items():
            texts = [it.text for it in group]
            original_texts = list(texts)
            context = make_context(sidx)

            if not should_translate_group(sidx):
                if log_fh:
                    log_fh.write(
                        f"## Group owner: {group[0].owner if group else 'unknown'}\n"
                    )
                if log_fh:
                    log_fh.write(f"Context:\n{context}\n")
                    log_fh.write("Skipped due to --slides filter.\n\n")
                progress.update(len(group))
                continue

            non_empty_idxs = [i for i, t in enumerate(texts) if t.strip()]
            if not non_empty_idxs:
                if log_fh:
                    log_fh.write(
                        f"## Group owner: {group[0].owner if group else 'unknown'}\n"
                    )
                    log_fh.write(f"Context:\n{context}\n")
                    log_fh.write("All items blank; skipped translation.\n\n")
                progress.update(len(group))
                continue

            to_send = [texts[i] for i in non_empty_idxs]
            translated = translator.translate(to_send, context=context)

            any_translated = any_translated or bool(translated)

            j = 0
            for i in range(len(texts)):
                if i in non_empty_idxs:
                    new_t = translated[j]
                    if glossary:
                        new_t = apply_glossary(new_t, glossary)
                    texts[i] = new_t
                    j += 1

            if log_fh:
                log_fh.write(
                    f"## Group owner: {group[0].owner if group else 'unknown'}\n"
                )
                log_fh.write(f"Context:\n{context}\n")
                for item, src_text, dst_text in zip(group, original_texts, texts):
                    log_fh.write(f"- {item.owner} (#{item.idx})\n")
                    log_fh.write("SRC:\n")
                    log_fh.write((src_text or "") + "\n")
                    log_fh.write("DST:\n")
                    log_fh.write((dst_text or "") + "\n\n")

            if dry_run:
                for item, new_t in zip(group, texts):
                    if new_t != item.text:
                        print(f"[DRY-RUN] {item.owner} (#{item.idx}):")
                        src = item.text
                        dst = new_t

                        def trunc(s: str) -> str:
                            return (s[:120] + "...") if len(s) > 120 else s

                        print("  SRC:", repr(trunc(src)))
                        print("  DST:", repr(trunc(dst)))
            for k, item in enumerate(group):
                items[item.idx] = TextItem(owner=item.owner, idx=item.idx, text=texts[k])

            progress.update(len(group))

        progress.close()

        if dry_run:
            print("Dry-run complete. No file written.")
            return TranslationResult(skipped_all=not any_translated, warnings=warnings)

        prs2 = Presentation(input_path)
        new_items, _ = collect_items(
            prs2, include_notes=include_notes, include_masters=include_masters
        )
        if len(new_items) != len(items):
            warnings.append(
                "[WARN] Item count changed between passes; attempting best-effort application."
            )

        idx = 0
        for sidx, slide in enumerate(prs2.slides):
            for tf in iter_text_frames(slide.shapes):
                if idx < len(items) and items[idx].owner.startswith("slide["):
                    set_text(tf, items[idx].text)
                    idx += 1
            if include_notes and slide.has_notes_slide:
                try:
                    ntf = slide.notes_slide.notes_text_frame
                    if ntf:
                        if idx < len(items) and items[idx].owner.startswith(
                            f"slide[{sidx+1}]-notes"
                        ):
                            set_text(ntf, items[idx].text)
                            idx += 1
                except Exception:
                    pass

        if include_masters:
            for m_idx, master in enumerate(prs2.slide_masters):
                for tf in iter_text_frames(master.shapes):
                    if idx < len(items) and items[idx].owner.startswith(
                        f"master[{m_idx+1}]"
                    ):
                        set_text(tf, items[idx].text)
                        idx += 1
                for l_idx, layout in enumerate(master.slide_layouts):
                    for tf in iter_text_frames(layout.shapes):
                        if idx < len(items) and items[idx].owner.startswith(
                            f"layout[{m_idx+1}.{l_idx+1}]"
                        ):
                            set_text(tf, items[idx].text)
                            idx += 1

        prs2.save(output_path)
        print(f"Saved translated presentation to: {output_path}")

        return TranslationResult(skipped_all=not any_translated, warnings=warnings)
    finally:
        if log_fh:
            log_fh.close()
