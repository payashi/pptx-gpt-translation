from dataclasses import dataclass
from typing import Dict, List, Optional, Set, Tuple

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.table import Table


@dataclass
class TextItem:
    owner: str  # e.g., "slide[3]" or "layout[1.2]" etc.
    idx: int  # stable index in collection
    text: str


def parse_slide_spec(spec: Optional[str], total_slides: int) -> Tuple[Optional[Set[int]], bool]:
    if not spec:
        return None, False
    selected: Set[int] = set()
    invalid = False
    parts = [p.strip() for p in spec.split(",") if p.strip()]
    for part in parts:
        if "-" in part:
            start_str, end_str = part.split("-", 1)
            if not start_str.strip().isdigit() or not end_str.strip().isdigit():
                invalid = True
                continue
            start = int(start_str)
            end = int(end_str)
            if end < start:
                start, end = end, start
            for num in range(start, end + 1):
                if 1 <= num <= total_slides:
                    selected.add(num - 1)
                else:
                    invalid = True
        else:
            if not part.isdigit():
                invalid = True
                continue
            num = int(part)
            if 1 <= num <= total_slides:
                selected.add(num - 1)
            else:
                invalid = True
    return selected, invalid


def iter_text_frames(shapes):
    # Yield text_frame objects from shapes (including groups, tables, and some charts)
    for shape in shapes:
        chart = None
        try:
            chart = shape.chart
        except (AttributeError, ValueError, KeyError):
            chart = None
        except Exception:
            chart = None

        if chart is not None:
            try:
                if (
                    chart.has_title
                    and chart.chart_title
                    and chart.chart_title.has_text_frame
                ):
                    yield chart.chart_title.text_frame
            except Exception:
                pass
            try:
                for axis_name in ("category_axis", "value_axis", "series_axis"):
                    axis = getattr(chart, axis_name, None)
                    if axis is not None and getattr(axis, "has_title", False):
                        at = getattr(axis, "axis_title", None)
                        if at is not None and getattr(at, "has_text_frame", False):
                            yield at.text_frame
            except Exception:
                pass

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for tf in iter_text_frames(shape.shapes):
                yield tf
        elif getattr(shape, "has_text_frame", False):
            if shape.has_text_frame and shape.text_frame is not None:
                yield shape.text_frame
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table: Table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    if cell.text_frame:
                        yield cell.text_frame


def get_text(tf) -> str:
    paras = []
    for p in tf.paragraphs:
        runs = [r.text for r in p.runs] or [p.text]
        paras.append("".join(runs))
    return "\n".join(paras).strip()


def set_text(tf, new_text: str):
    # Preserve existing formatting by reusing current paragraphs and runs.
    new_text = "" if new_text is None else new_text
    paragraphs = list(tf.paragraphs)

    # If there is no formatting to preserve, fall back to direct assignment.
    if not paragraphs:
        tf.text = new_text
        return

    parts = new_text.split("\n")

    if len(parts) <= len(paragraphs):
        assign_count = len(parts)
    else:
        assign_count = len(paragraphs)

    for idx in range(assign_count):
        text_piece = parts[idx]
        if idx == assign_count - 1 and len(parts) > len(paragraphs):
            overflow = parts[idx:]
            text_piece = "\n".join(overflow)
        para = paragraphs[idx]
        if para.runs:
            primary_run = para.runs[0]
        else:
            primary_run = para.add_run()
        primary_run.text = text_piece
        for extra in para.runs[1:]:
            extra.text = ""

    # Remove trailing paragraphs when translation has fewer segments.
    for para in reversed(paragraphs[assign_count:]):
        p = para._p
        parent = p.getparent()
        if parent is not None:
            parent.remove(p)


def collect_items(prs: Presentation, include_notes: bool, include_masters: bool):
    items: List[TextItem] = []
    slide_titles: List[str] = []
    # Slides
    for sidx, slide in enumerate(prs.slides):
        # Try to extract a plausible slide title (first placeholder title or first text box)
        title_text = ""
        try:
            if slide.shapes.title and getattr(
                slide.shapes.title, "has_text_frame", False
            ):
                title_text = get_text(slide.shapes.title.text_frame)
        except Exception:
            pass
        if not title_text:
            # fallback: first text frame on slide
            for tf in iter_text_frames(slide.shapes):
                title_text = get_text(tf)
                if title_text:
                    break
        slide_titles.append(title_text[:120])

        for tf in iter_text_frames(slide.shapes):
            items.append(
                TextItem(owner=f"slide[{sidx+1}]", idx=len(items), text=get_text(tf))
            )

        if include_notes and slide.has_notes_slide:
            try:
                ntf = slide.notes_slide.notes_text_frame
                if ntf:
                    items.append(
                        TextItem(
                            owner=f"slide[{sidx+1}]-notes",
                            idx=len(items),
                            text=get_text(ntf),
                        )
                    )
            except Exception:
                pass

    # Masters/layouts
    if include_masters:
        for m_idx, master in enumerate(prs.slide_masters):
            for tf in iter_text_frames(master.shapes):
                items.append(
                    TextItem(
                        owner=f"master[{m_idx+1}]", idx=len(items), text=get_text(tf)
                    )
                )
            for l_idx, layout in enumerate(master.slide_layouts):
                for tf in iter_text_frames(layout.shapes):
                    items.append(
                        TextItem(
                            owner=f"layout[{m_idx+1}.{l_idx+1}]",
                            idx=len(items),
                            text=get_text(tf),
                        )
                    )
    return items, slide_titles


def summarize_deck(slide_titles: List[str], max_slides: int = 4) -> str:
    if not slide_titles:
        return ""
    first = slide_titles[:max_slides]
    summary = "Deck overview titles:\n- " + "\n- ".join([t for t in first if t])
    return summary
