#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Translate PPTX with GPT + python-pptx.

- Batches by slide and uses neighbor/summary context to improve coherence.
- Optional glossary enforcement (post-pass string replacement).
- Optional speaker notes and masters/layouts.

Requires:
  pip install openai python-pptx tqdm tenacity
  export OPENAI_API_KEY=...

Example:
  python gpt_pptx_translate.py input.pptx output_ja.pptx --target JA --source EN --notes --masters --glossary glossary.csv --model gpt-4o-mini --strategy neighbor
"""
import os
import csv
import argparse
from typing import List, Tuple, Optional, Dict, Set
from dataclasses import dataclass

from tenacity import (
    retry,
    stop_after_attempt,
    wait_exponential,
    retry_if_exception_type,
)
from pptx import Presentation
from pptx.table import Table
from pptx.enum.shapes import MSO_SHAPE_TYPE
from tqdm import tqdm
from dotenv import load_dotenv

load_dotenv()

# OpenAI SDK v1.x
try:
    from openai import OpenAI
except Exception as e:
    OpenAI = None  # handled in main

# -----------------------------
# Utilities
# -----------------------------


@dataclass
class TextItem:
    owner: str  # e.g., "slide[3]" or "layout[1.2]" etc.
    idx: int  # stable index in collection
    text: str


def read_glossary(path: Optional[str]) -> List[Tuple[str, str]]:
    pairs: List[Tuple[str, str]] = []
    if not path:
        return pairs
    with open(path, newline="", encoding="utf-8-sig") as f:
        reader = csv.reader(f)
        for row in reader:
            if len(row) >= 2:
                src, tgt = row[0].strip(), row[1].strip()
                if src:
                    pairs.append((src, tgt))
    return pairs


def apply_glossary(text: str, glossary: List[Tuple[str, str]]) -> str:
    for src, tgt in glossary:
        if src:
            text = text.replace(src, tgt)
    return text


def parse_slide_spec(spec: Optional[str], total_slides: int) -> Tuple[Optional[Set[int]], bool]:
    if not spec:
        return None, False
    selected: Set[int] = set()
    invalid = False
    parts = [p.strip() for p in spec.split(",") if p.strip()]
    for part in parts:
        if "-" in part:
            start_str, end_str = part.split("-", 1)
            if not start_str.strip().isdigit() or not end_str.strip().isdigit():
                invalid = True
                continue
            start = int(start_str)
            end = int(end_str)
            if end < start:
                start, end = end, start
            for num in range(start, end + 1):
                if 1 <= num <= total_slides:
                    selected.add(num - 1)
                else:
                    invalid = True
        else:
            if not part.isdigit():
                invalid = True
                continue
            num = int(part)
            if 1 <= num <= total_slides:
                selected.add(num - 1)
            else:
                invalid = True
    return selected, invalid


def iter_text_frames(shapes):
    # Yield text_frame objects from shapes (including groups, tables, and some charts)
    for shape in shapes:
        # Charts (where available)
        chart = None
        try:
            chart = shape.chart
        except (AttributeError, ValueError, KeyError):
            chart = None
        except Exception:
            chart = None

        if chart is not None:
            try:
                if (
                    chart.has_title
                    and chart.chart_title
                    and chart.chart_title.has_text_frame
                ):
                    yield chart.chart_title.text_frame
            except Exception:
                pass
            try:
                for axis_name in ("category_axis", "value_axis", "series_axis"):
                    axis = getattr(chart, axis_name, None)
                    if axis is not None and getattr(axis, "has_title", False):
                        at = getattr(axis, "axis_title", None)
                        if at is not None and getattr(at, "has_text_frame", False):
                            yield at.text_frame
            except Exception:
                pass

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for tf in iter_text_frames(shape.shapes):
                yield tf
        elif getattr(shape, "has_text_frame", False):
            if shape.has_text_frame and shape.text_frame is not None:
                yield shape.text_frame
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table: Table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    if cell.text_frame:
                        yield cell.text_frame


def get_text(tf) -> str:
    paras = []
    for p in tf.paragraphs:
        runs = [r.text for r in p.runs] or [p.text]
        paras.append("".join(runs))
    return "\n".join(paras).strip()


def set_text(tf, new_text: str):
    # Preserve existing formatting by reusing current paragraphs and runs.
    new_text = "" if new_text is None else new_text
    paragraphs = list(tf.paragraphs)

    # If there is no formatting to preserve, fall back to direct assignment.
    if not paragraphs:
        tf.text = new_text
        return

    parts = new_text.split("\n")

    if len(parts) <= len(paragraphs):
        assign_count = len(parts)
    else:
        assign_count = len(paragraphs)

    for idx in range(assign_count):
        text_piece = parts[idx]
        if idx == assign_count - 1 and len(parts) > len(paragraphs):
            overflow = parts[idx:]
            text_piece = "\n".join(overflow)
        para = paragraphs[idx]
        if para.runs:
            primary_run = para.runs[0]
        else:
            primary_run = para.add_run()
        primary_run.text = text_piece
        for extra in para.runs[1:]:
            extra.text = ""

    # Remove trailing paragraphs when translation has fewer segments.
    for para in reversed(paragraphs[assign_count:]):
        p = para._p
        parent = p.getparent()
        if parent is not None:
            parent.remove(p)


def collect_items(prs: Presentation, include_notes: bool, include_masters: bool):
    items: List[TextItem] = []
    slide_titles: List[str] = []
    # Slides
    for sidx, slide in enumerate(prs.slides):
        # Try to extract a plausible slide title (first placeholder title or first text box)
        title_text = ""
        try:
            if slide.shapes.title and getattr(
                slide.shapes.title, "has_text_frame", False
            ):
                title_text = get_text(slide.shapes.title.text_frame)
        except Exception:
            pass
        if not title_text:
            # fallback: first text frame on slide
            for tf in iter_text_frames(slide.shapes):
                title_text = get_text(tf)
                if title_text:
                    break
        slide_titles.append(title_text[:120])

        for tf in iter_text_frames(slide.shapes):
            items.append(
                TextItem(owner=f"slide[{sidx+1}]", idx=len(items), text=get_text(tf))
            )

        if include_notes and slide.has_notes_slide:
            try:
                ntf = slide.notes_slide.notes_text_frame
                if ntf:
                    items.append(
                        TextItem(
                            owner=f"slide[{sidx+1}]-notes",
                            idx=len(items),
                            text=get_text(ntf),
                        )
                    )
            except Exception:
                pass

    # Masters/layouts
    if include_masters:
        for m_idx, master in enumerate(prs.slide_masters):
            for tf in iter_text_frames(master.shapes):
                items.append(
                    TextItem(
                        owner=f"master[{m_idx+1}]", idx=len(items), text=get_text(tf)
                    )
                )
            for l_idx, layout in enumerate(master.slide_layouts):
                for tf in iter_text_frames(layout.shapes):
                    items.append(
                        TextItem(
                            owner=f"layout[{m_idx+1}.{l_idx+1}]",
                            idx=len(items),
                            text=get_text(tf),
                        )
                    )
    return items, slide_titles


def summarize_deck(slide_titles: List[str], max_slides: int = 4) -> str:
    if not slide_titles:
        return ""
    first = slide_titles[:max_slides]
    summary = "Deck overview titles:\n- " + "\n- ".join([t for t in first if t])
    return summary


# -----------------------------
# GPT Translator
# -----------------------------


class GPTTranslator:
    def __init__(
        self,
        model: str,
        source: Optional[str],
        target: str,
        temperature: Optional[float] = None,
    ):
        if OpenAI is None:
            raise RuntimeError("OpenAI SDK not available. Install `openai` >= 1.37.0.")
        api_key = os.environ.get("OPENAI_API_KEY")
        if not api_key:
            raise RuntimeError("Set OPENAI_API_KEY environment variable.")
        self.client = OpenAI(api_key=api_key)
        self.model = model
        self.source = source
        self.target = target
        self.temperature = temperature

    @retry(
        reraise=True,
        stop=stop_after_attempt(4),
        wait=wait_exponential(multiplier=1, min=1, max=15),
        retry=retry_if_exception_type(Exception),
    )
    def translate(self, texts: List[str], context: str = "") -> List[str]:
        """
        Translate a list of texts as a block, returning a same-length list.
        Uses a delimiter to split/merge. We keep it simple to minimize JSON parsing issues.
        """
        if not texts:
            return []

        delimiter = "\n\n<<<SPLIT>>>\n\n"
        joined = delimiter.join(texts)

        sys = (
            "You are a professional translator. "
            f"Translate from {self.source or 'the source language'} to {self.target}. "
            "Preserve technical terms, numbers, math, and code blocks. "
            "Keep bullet-like brevity for short lines; keep paragraph flow for long text. "
            "While faithfully retaining every keyword from the source (device names, terminology, etc.), craft the translation so it stays natural and slide-ready: concise and preferably ending in nouns. "
            "If a line is a source attribution, copy it verbatim without translating. "
            "Do NOT add extra commentary. Return only the translations joined by the same delimiter."
        )

        if context:
            sys += f" Use this context for disambiguation: {context[:4000]}"

        request_args = {
            "model": self.model,
            "messages": [
                {"role": "system", "content": sys},
                {"role": "user", "content": joined},
            ],
        }
        if self.temperature is not None:
            request_args["temperature"] = self.temperature

        resp = self.client.chat.completions.create(**request_args)
        out = resp.choices[0].message.content or ""
        parts = out.split(delimiter)
        # If counts mismatch, fall back to naive split by lines (rare)
        if len(parts) != len(texts):
            # Attempt a softer split: keep lengths equal
            # We'll just distribute by paragraph counts as a fallback
            return out.splitlines()[: len(texts)] + [""] * (
                len(texts) - len(out.splitlines())
            )
        return parts


# -----------------------------
# Main CLI
# -----------------------------


def main():
    ap = argparse.ArgumentParser(
        description="Translate PPTX slides with GPT + python-pptx."
    )
    ap.add_argument("input", help="Input .pptx path")
    ap.add_argument("output", help="Output .pptx path")
    ap.add_argument(
        "--model",
        default="gpt-4o-mini",
        help="OpenAI model (e.g., gpt-4o, gpt-4o-mini)",
    )
    ap.add_argument("--source", default=None, help="Source language (e.g., EN)")
    ap.add_argument("--target", required=True, help="Target language (e.g., JA)")
    ap.add_argument("--notes", action="store_true", help="Include speaker notes")
    ap.add_argument(
        "--masters", action="store_true", help="Include slide masters/layouts"
    )
    ap.add_argument(
        "--glossary", default=None, help="CSV file with source,target terms"
    )
    ap.add_argument(
        "--strategy",
        default="neighbor",
        choices=["neighbor", "title-only", "deck"],
        help="Context strategy",
    )
    ap.add_argument(
        "--temperature",
        type=float,
        default=None,
        help="Sampling temperature (omit to use the model's default)",
    )
    ap.add_argument(
        "--dry_run", action="store_true", help="Preview changes without saving"
    )
    ap.add_argument("--log", default=None, help="File path to write translation log")
    ap.add_argument(
        "--slides",
        default=None,
        help="Comma-separated slide numbers or ranges to translate (e.g., 1,3-5)",
    )
    args = ap.parse_args()

    prs = Presentation(args.input)
    items, slide_titles = collect_items(
        prs, include_notes=args.notes, include_masters=args.masters
    )
    selected_slides, slides_invalid = parse_slide_spec(args.slides, len(prs.slides))
    if slides_invalid:
        print("[WARN] Some slide identifiers in --slides were invalid or out of range and have been ignored.")
    if args.slides and selected_slides is not None and not selected_slides:
        print("[WARN] No valid slide numbers specified in --slides; nothing will be translated.")

    # Map slide index by owner string
    slide_indices = []
    for it in items:
        # Extract slide index if present like slide[3] or slide[3]-notes
        sidx = None
        if it.owner.startswith("slide["):
            try:
                sidx = int(it.owner.split("[")[1].split("]")[0]) - 1
            except Exception:
                sidx = None
        slide_indices.append(sidx)

    glossary = read_glossary(args.glossary)
    translator = GPTTranslator(
        model=args.model,
        source=args.source,
        target=args.target,
        temperature=args.temperature,
    )

    log_fh = None
    if args.log:
        log_fh = open(args.log, "w", encoding="utf-8")

    try:
        # Group items by slide index for batching (masters/layouts may be None; we still batch them separately)
        by_slide: Dict[Optional[int], List[TextItem]] = {}
        for it, sidx in zip(items, slide_indices):
            by_slide.setdefault(sidx, []).append(it)

        # Build deck summary once
        deck_summary = summarize_deck(slide_titles)

        # Translate per group
        progress = tqdm(total=len(items), desc="Translating")

        # Helper to make context per slide
        def make_context(current_idx: Optional[int]) -> str:
            if args.strategy == "title-only":
                title = (
                    slide_titles[current_idx]
                    if (current_idx is not None and 0 <= current_idx < len(slide_titles))
                    else ""
                )
                return f"Deck title(s): {slide_titles[0] if slide_titles else ''}. Current slide title: {title}."
            elif args.strategy == "deck":
                return deck_summary
            else:  # neighbor
                if current_idx is None:
                    return deck_summary
                titles = []
                for j in [current_idx - 1, current_idx, current_idx + 1]:
                    if 0 <= j < len(slide_titles):
                        titles.append(f"Slide {j+1}: {slide_titles[j]}")
                return "Neighbor titles:\n" + "\n".join(titles)

        def should_translate_group(group_slide_idx: Optional[int]) -> bool:
            if selected_slides is None:
                return True
            if group_slide_idx is None:
                return False
            return group_slide_idx in selected_slides

        # Perform translation
        for sidx, group in by_slide.items():
            texts = [it.text for it in group]
            original_texts = list(texts)
            context = make_context(sidx)
            if not should_translate_group(sidx):
                if log_fh:
                    log_fh.write(f"## Group owner: {group[0].owner if group else 'unknown'}\n")
                    log_fh.write(f"Context:\n{context}\n")
                    log_fh.write("Skipped due to --slides filter.\n\n")
                progress.update(len(group))
                continue
            # Skip empty batch quickly
            non_empty_idxs = [i for i, t in enumerate(texts) if t.strip()]
            if not non_empty_idxs:
                if log_fh:
                    log_fh.write(f"## Group owner: {group[0].owner if group else 'unknown'}\n")
                    log_fh.write(f"Context:\n{context}\n")
                    log_fh.write("All items blank; skipped translation.\n\n")
                progress.update(len(group))
                continue
            to_send = [texts[i] for i in non_empty_idxs]
            translated = translator.translate(to_send, context=context)
            # Place back
            j = 0
            for i in range(len(texts)):
                if i in non_empty_idxs:
                    new_t = translated[j]
                    if glossary:
                        new_t = apply_glossary(new_t, glossary)
                    texts[i] = new_t
                    j += 1
            if log_fh:
                log_fh.write(f"## Group owner: {group[0].owner if group else 'unknown'}\n")
                log_fh.write(f"Context:\n{context}\n")
                for item, src_text, dst_text in zip(group, original_texts, texts):
                    log_fh.write(f"- {item.owner} (#{item.idx})\n")
                    log_fh.write("SRC:\n")
                    log_fh.write((src_text or "") + "\n")
                    log_fh.write("DST:\n")
                    log_fh.write((dst_text or "") + "\n\n")
            # Write back or preview
            if args.dry_run:
                for item, new_t in zip(group, texts):
                    if new_t != item.text:
                        print(f"[DRY-RUN] {item.owner} (#{item.idx}):")
                        src = item.text
                        dst = new_t

                        def trunc(s):
                            return (s[:120] + "...") if len(s) > 120 else s

                        print("  SRC:", repr(trunc(src)))
                        print("  DST:", repr(trunc(dst)))
            else:
                # Re-apply into the actual text frames
                # We need to find the shape again; simplest approach: iterate again in the same order
                # We'll re-collect frames and set in a second pass to keep the mapping stable.
                pass  # defer to a final application loop after translation
            # Store final texts back into items array for later application
            for k, item in enumerate(group):
                items[item.idx] = TextItem(owner=item.owner, idx=item.idx, text=texts[k])

            progress.update(len(group))

        progress.close()

        if args.dry_run:
            print("Dry-run complete. No file written.")
            return

        # Final pass: re-open and write translated strings in the same discovery order
        prs2 = Presentation(args.input)
        new_items, _ = collect_items(
            prs2, include_notes=args.notes, include_masters=args.masters
        )
        if len(new_items) != len(items):
            print(
                "[WARN] Item count changed between passes; attempting best-effort application."
            )
        apply_count = 0
        for old, (owner, new) in zip(items, [(ni.owner, ni.text) for ni in items]):
            # We must step through new_items in the same sequence and set text
            # We rely on the order of traversal being stable for the same file.
            pass

        # Since we can't directly store references to text_frames across instances,
        # repeat traversal and set in place:
        idx = 0
        for sidx, slide in enumerate(prs2.slides):
            # slide shapes
            for tf in iter_text_frames(slide.shapes):
                if idx < len(items) and items[idx].owner.startswith("slide["):
                    set_text(tf, items[idx].text)
                    idx += 1
            # notes
            if args.notes and slide.has_notes_slide:
                try:
                    ntf = slide.notes_slide.notes_text_frame
                    if ntf:
                        if idx < len(items) and items[idx].owner.startswith(
                            f"slide[{sidx+1}]-notes"
                        ):
                            set_text(ntf, items[idx].text)
                            idx += 1
                except Exception:
                    pass

        # Masters/layouts
        if args.masters:
            for m_idx, master in enumerate(prs2.slide_masters):
                for tf in iter_text_frames(master.shapes):
                    if idx < len(items) and items[idx].owner.startswith(
                        f"master[{m_idx+1}]"
                    ):
                        set_text(tf, items[idx].text)
                        idx += 1
                for l_idx, layout in enumerate(master.slide_layouts):
                    for tf in iter_text_frames(layout.shapes):
                        if idx < len(items) and items[idx].owner.startswith(
                            f"layout[{m_idx+1}.{l_idx+1}]"
                        ):
                            set_text(tf, items[idx].text)
                            idx += 1

        prs2.save(args.output)
        print(f"Saved translated presentation to: {args.output}")
    finally:
        if log_fh:
            log_fh.close()


if __name__ == "__main__":
    main()
